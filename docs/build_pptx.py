#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
LHAMS 세미나 PPT 빌드 스크립트
================================
기존 docs/LHAMS_크리니티_통합세미나.pptx(20슬라이드)를 열어:
  1) 최신 기능(관리 API/감사이력/체크리스트/요약탭/SI·폐쇄망 대응)을 반영해
     일부 슬라이드의 텍스트만 갱신하고
  2) "소스별 역할" 매핑 + "관리 기능 고도화" + "요약·SI 표준화" 4개 슬라이드를
     기존 디자인 톤(다크 테마, Consolas/맑은고딕, 티얼 포인트)에 맞춰 새로 추가한 뒤
  3) 전체 페이지 번호·섹션 번호를 재계산해 같은 파일명으로 저장한다.

재실행 가능하도록 항상 "원본 스타일 상수"만 사용해 새 슬라이드를 그리며,
기존 슬라이드는 지오메트리를 건드리지 않고 텍스트만 치환한다(레이아웃 손상 위험 최소화).
"""

import copy
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

DOCS_DIR = os.path.dirname(os.path.abspath(__file__))
PPTX_PATH = os.path.join(DOCS_DIR, "LHAMS_크리니티_통합세미나.pptx")
DASHBOARD_IMG = os.path.join(DOCS_DIR, "assets", "dashboard_live.png")

# ─────────────────────────── 디자인 토큰 (기존 슬라이드에서 추출) ───────────────────────────
BG = RGBColor(0x0E, 0x14, 0x1B)
PANEL = RGBColor(0x15, 0x1D, 0x26)
PANEL2 = RGBColor(0x1A, 0x24, 0x30)
LINE = RGBColor(0x24, 0x30, 0x3C)
TEXT = RGBColor(0xC9, 0xD4, 0xDF)
MUTED = RGBColor(0x6E, 0x7E, 0x8E)
WHITE_ISH = RGBColor(0xF4, 0xF8, 0xFB)
CLEAN = RGBColor(0x3D, 0xB8, 0x8B)
WARN = RGBColor(0xF0, 0xA6, 0x3C)
HIGH = RGBColor(0xE8, 0x73, 0x4D)
CRIT = RGBColor(0xE8, 0x55, 0x4D)
ACCENT = RGBColor(0x7F, 0xB3, 0xD5)

MONO = "Consolas"
SANS = "맑은 고딕"

MARGIN_L = 0.5
CONTENT_R = 12.83  # 마진 우측 끝 (슬라이드 폭 13.333in 기준)


# ─────────────────────────── 저수준 헬퍼 ───────────────────────────
def find_layout(prs, name="Blank"):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    return prs.slide_layouts[6]


def new_slide(prs, layout):
    slide = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def _set_run(run, text, font=SANS, size=14, bold=False, color=TEXT):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def textbox(slide, left, top, width, height, lines, font=SANS, size=14, bold=False,
            color=TEXT, align=PP_ALIGN.LEFT, word_wrap=True):
    """lines: str 한 줄이면 그대로, 리스트면 여러 단락으로."""
    if isinstance(lines, str):
        lines = [lines]
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = word_wrap
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        _set_run(p.add_run(), line, font, size, bold, color)
    return box


def header(slide, top_right_label, page_no, total):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(0.34), Inches(0.09), Inches(0.09))
    dot.fill.solid(); dot.fill.fore_color.rgb = CLEAN; dot.line.fill.background(); dot.shadow.inherit = False

    textbox(slide, 0.68, 0.28, 2.6, 0.3, "LHAMS", font=MONO, size=13, bold=True, color=WHITE_ISH)
    textbox(slide, 1.55, 0.29, 4.0, 0.28, "보안 감사 시스템 세미나", font=SANS, size=10.5, color=MUTED)
    textbox(slide, 9.6, 0.29, 3.2, 0.28, top_right_label, font=MONO, size=10, color=MUTED, align=PP_ALIGN.RIGHT)

    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.5), Inches(0.72), Inches(CONTENT_R), Inches(0.72))
    conn.line.color.rgb = LINE
    conn.line.width = Pt(0.75)

    textbox(slide, 11.9, 7.08, 0.9, 0.3, f"{page_no:02d} / {total}", font=MONO, size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def section_tag(slide, tag_text):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.01), Inches(0.22), Inches(0.019))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background(); bar.shadow.inherit = False
    textbox(slide, 0.8, 0.92, 8.0, 0.3, tag_text, font=MONO, size=11.5, bold=True, color=ACCENT)


def headline(slide, text, top=1.28, size=28):
    textbox(slide, 0.5, top, 12.3, 0.9, text, font=SANS, size=size, bold=True, color=WHITE_ISH)


def body_slide(prs, layout, top_right_label, page_no, total, tag, title):
    slide = new_slide(prs, layout)
    header(slide, top_right_label, page_no, total)
    section_tag(slide, tag)
    headline(slide, title)
    return slide


def impact_box(slide, text, top=4.75, height=1.15, label="IMPACT", label_color=ACCENT):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(top), Inches(11.8), Inches(height))
    box.fill.solid(); box.fill.fore_color.rgb = PANEL2; box.line.fill.background(); box.shadow.inherit = False
    textbox(slide, 0.75, top + 0.13, 2.2, 0.3, label, font=MONO, size=10.5, bold=True, color=label_color)
    textbox(slide, 0.75, top + 0.45, 11.3, height - 0.55, text, font=SANS, size=13, color=TEXT)


def bullets(slide, lines, top=1.95, width=11.8, size=14, line_h=0.5):
    textbox(slide, 0.5, top, width, line_h * len(lines) + 0.2, lines, font=SANS, size=size, color=TEXT)


def role_table(slide, rows, top=1.55, name_w=3.3, desc_w=8.5, row_h=0.46, name_size=12, desc_size=11.5, header=None):
    name_x = 0.5
    desc_x = name_x + name_w + 0.15
    y = top
    if header:
        textbox(slide, name_x, y, name_w, 0.3, header[0], font=MONO, size=10, bold=True, color=MUTED)
        textbox(slide, desc_x, y, desc_w, 0.3, header[1], font=MONO, size=10, bold=True, color=MUTED)
        y += 0.34
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(name_x), Inches(y), Inches(desc_x + desc_w), Inches(y))
        line.line.color.rgb = LINE; line.line.width = Pt(0.75)
        y += 0.12
    for name, desc in rows:
        textbox(slide, name_x, y, name_w, row_h, name, font=MONO, size=name_size, bold=True, color=ACCENT)
        textbox(slide, desc_x, y, desc_w, row_h, desc, font=SANS, size=desc_size, color=TEXT)
        y += row_h
    return y


# ─────────────────────────── 텍스트 치환 헬퍼 (기존 슬라이드 편집용) ───────────────────────────
def find_shape_by_name(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    raise KeyError(f"shape not found: {name}")


def set_textbox_lines(shape, lines, font=None, size=None, bold=None, color=None):
    """기존 첫 run의 서식을 재사용하며 텍스트만 교체 (font/size/bold/color 지정 시 덮어씀)."""
    tf = shape.text_frame
    base = tf.paragraphs[0].runs[0]
    f_font = font or base.font.name
    f_size = size or (base.font.size.pt if base.font.size else 14)
    f_bold = base.font.bold if bold is None else bold
    try:
        f_color = color or base.font.color.rgb
    except Exception:
        f_color = color or TEXT
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_run(p.add_run(), line, f_font, f_size, f_bold, f_color)


def replace_by_prefix(slide, prefix, new_lines):
    """text_frame 전체 텍스트가 prefix로 시작하는 첫 shape을 찾아 교체."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip().startswith(prefix):
            set_textbox_lines(shape, new_lines)
            return True
    return False


def update_header_footer(slide, top_right_label=None, tag=None, page_no=None, total=None):
    if top_right_label is not None:
        set_textbox_lines(find_shape_by_name(slide, "TextBox 4"), [top_right_label])
    if tag is not None:
        set_textbox_lines(find_shape_by_name(slide, "TextBox 8"), [tag])
    if page_no is not None:
        set_textbox_lines(find_shape_by_name(slide, "TextBox 6"), [f"{page_no:02d} / {total}"])


# ─────────────────────────── 본 스크립트 ───────────────────────────
def main():
    prs = Presentation(PPTX_PATH)
    blank = find_layout(prs)

    slides = prs.slides

    # ---- 1) 기존 슬라이드 콘텐츠 갱신 (지오메트리는 그대로, 텍스트만 교체) ----

    # 03-3 · MIDDLEWARE (index 6, physical 7) — Flask 관리 API·자동 백업 언급 추가
    s = slides[6]
    replace_by_prefix(s, "—  각 센서가 남긴 텍스트 로그를 실시간 파싱", [
        "—  각 센서가 남긴 텍스트 로그를 실시간 파싱해 정형 JSON(lhams_audit.json)으로 적재",
        "—  TimedRotatingFileHandler로 자정마다 로그 자동 로테이션 + Flask 관리 API 백그라운드 기동",
        "—  경로/무시규칙/체크리스트/감사이력/헬스체크(/api/health)를 재시작 없이 운영, 변경 시 자동 백업",
    ])

    # 03-4 · FRONTEND (index 7, physical 8) — 3-tab 구조 언급
    s = slides[7]
    replace_by_prefix(s, "—  JSON 로그를 3초 주기로 폴링", [
        "—  JSON 로그 3초 폴링 + 관리 API 5초 헬스체크 · 요약/대시보드/관리자 3-tab 구조로 확장",
    ])

    # 07 · DASHBOARD (index 13, physical 14) — 캡션 갱신 + 오버사이즈 스크린샷 교정
    s = slides[13]
    replace_by_prefix(s, "React + Tailwind/SCSS", [
        'React + Vite + SCSS, 3초 폴링 · "심야 관제실" 다크 테마 — 대시보드는 요약·대시보드·관리자 '
        "3-tab 중 하나이며, 아래는 대시보드 탭 실제 화면입니다.",
    ])
    # 기존 파일의 Picture/테두리 Rectangle이 슬라이드 하단을 벗어나던 문제 교정 (비율 유지 축소)
    pic_left, pic_top, pic_w, pic_h = 3.47, 2.3, 6.4, 4.5
    for shape in s.shapes:
        if shape.name in ("Picture 11", "Rectangle 12"):
            shape.left, shape.top, shape.width, shape.height = (
                Inches(pic_left), Inches(pic_top), Inches(pic_w), Inches(pic_h)
            )

    # 08 · BACKEND (index 14, physical 15) — 신규 영속 파일 언급 추가
    s = slides[14]
    replace_by_prefix(s, "—  on_created / on_modified", [
        "—  on_created / on_modified / on_deleted / on_moved 훅으로 전 이벤트 포착",
        "—  생성·수정 시 scan_malware()가 clamdscan 호출 → 탐지 시 risk_level 자동 격상",
        "—  ausearch 연동으로 파일 소유자(owner)와 별개로 실제 행위자(auid)까지 함께 기록",
        "—  최신 200건만 유지(MAX_JSON_EVENTS) — admin_audit.json(관리자 조작 감사)·checklist.json(구축 체크리스트)도 별도 관리",
    ])

    # 10 · PROJECT LAYOUT (index 16, physical 17) — 폴더 주석 갱신
    s = slides[16]
    tree_shape = find_shape_by_name(s, "TextBox 12")
    tf = tree_shape.text_frame
    rows = [
        ("docs/", "        # 세미나 · 아키텍처 문서"),
        ("frontend/", "    # React 관제 대시보드 (요약·대시보드·관리자 3-tab)"),
        ("agent/", "       # Python Watchdog 컨트롤 타워 + 관리 API(Flask)"),
        ("java/", "        # Tomcat 임베딩용 대안 에이전트 (FileAuditWatcher)"),
        ("scripts/", "     # inotify / ClamAV / auditd 자동화"),
        ("systemd/", "     # 자가 치유 서비스 유닛"),
        ("data/", "        # test_monitor · quarantine · logs · config_backups (Git 제외)"),
    ]
    tf.clear()
    for i, (name, comment) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_run(p.add_run(), name, MONO, 13, True, ACCENT)
        _set_run(p.add_run(), comment, MONO, 13, False, MUTED)

    # 11 · LIVE DEMO (index 17, physical 18 → 새 물리적 위치 22, 섹션 15) — EICAR 로컬화
    s = slides[17]
    replace_by_prefix(s, "touch → echo(수정)", [
        "touch → echo(수정) → mv(교체) → rm(삭제) → EICAR 로컬 생성(악성코드)",
    ])

    # 12 · CONCLUSION (index 18, physical 19 → 새 물리적 위치 23, 섹션 16) — 성과/로드맵 갱신
    s = slides[18]
    replace_by_prefix(s, "—  보안 위협에 대한 빠른 초기 대응", [
        "—  보안 위협에 대한 빠른 초기 대응 + 관리자 조작까지 포함한 이중 감사 체계",
        "—  서버 관리 리소스 절감 — 설정 백업/내보내기·가져오기로 다중 서버 표준화",
        "—  시스템 안정성 확보 + 헬스체크(/api/health)로 장애 조기 발견",
    ])
    replace_by_prefix(s, "—  취약점 진단 및 리포팅 체계화", [
        "—  계정 기반 로그인/SSO 연동으로 관리자 인증 고도화",
        "—  이상 징후 발생 시 Slack/이메일 자동 알림 연동",
        "—  다중 사이트 통합 모니터링 대시보드(전사 현황판)로 확장",
    ])

    # ---- 2) 신규 슬라이드 4개 추가 (마지막에 append 후 순서 이동) ----
    TOTAL = 24

    # 11-1 · SOURCE MAP (BACKEND) — physical 18
    s = body_slide(prs, blank, "11 — 소스맵·백엔드", 18, TOTAL, "11-1 · SOURCE MAP (BACKEND)",
                   "어떤 코드가 무엇을 하는가 — Python 백엔드")
    role_table(s, [
        ("ConfigStore", "data/config.json 로드·저장, 저장할 때마다 config_backups/에 자동 백업(최근 20개 보관)"),
        ("WatchManager", "감시 경로 등록·on/off·재귀 스케줄링, 실제 감시 성공 여부(active)·실패 사유(error) 추적"),
        ("QuarantineStore", "격리 파일의 원본 경로·시각·사유 기록, 복원·영구삭제 처리"),
        ("ChecklistStore", "신규 서버 구축 체크리스트 진행 상태(완료자·완료 시각) 관리"),
        ("AdminAuditLog", '관리자 조작(경로/설정/격리/체크리스트 변경)을 "누가·언제·무엇을" 기준으로 기록'),
        ("AuditdResolver", "ausearch로 auid(로그인 사용자)·comm(프로세스) 조회 — 파일 소유자와 별개로 실행 주체 특정"),
        ("LhamsAuditHandler", "watchdog 이벤트 훅 + clamdscan 검사 + JSON 적재 — 파일 감사의 핵심 엔진"),
        ("create_api() [Flask]", "/api/config·paths·quarantine·settings·checklist·admin-audit·health·config/export|import 라우팅"),
    ], top=1.9, row_h=0.42, header=("구성 요소", "역할"))
    impact_box(s, "파일 하나가 감시·검사·관리 API·감사·백업까지 담당하는 대신, 역할별 클래스로 분리해 "
                  "유지보수와 신규 기능 추가가 쉬워졌습니다.", top=5.85, height=0.9)

    # 11-2 · SOURCE MAP (FRONTEND) — physical 19
    s = body_slide(prs, blank, "11 — 소스맵·프론트", 19, TOTAL, "11-2 · SOURCE MAP (FRONTEND)",
                   "어떤 코드가 무엇을 하는가 — React 프론트엔드")
    role_table(s, [
        ("App.jsx", "이벤트 3초 폴링 + 시스템 진단(health) 5초 폴링, 요약/대시보드/관리자 탭 라우팅"),
        ("api.js", "관리 API 래퍼 — 변경 요청마다 작업자(actor)를 자동으로 실어 보냄"),
        ("SummaryPanel.jsx", "비전문가용 상태 배너·진단 카드·문장형 주요 이벤트 요약"),
        ("EventTable / StatCards", "이벤트 원장 테이블, 파일 교체(교체 전/후 경로) 강조 표시"),
        ("admin/ActorGate.jsx", '작업자 이름 식별 — 관리자 변경 이력의 "누가"를 채우는 입력 게이트'),
        ("admin/WatchPathManager", "감시 경로 등록/삭제/on-off + 경로 오류(⚠) 배지"),
        ("admin/QuarantineManager", "격리 파일 목록 확인, 복원/영구삭제"),
        ("admin/SetupChecklist", "설치 체크리스트 UI + 진행률 바"),
        ("admin/ConfigBackup", "설정 내보내기/가져오기(JSON) — 다중 서버 표준 배포"),
        ("admin/AdminAuditLog", "관리자 변경 이력 타임라인 UI"),
    ], top=1.9, row_h=0.36, name_size=11, desc_size=10.5, header=("구성 요소", "역할"))
    impact_box(s, "관리자 탭이 '보기 전용 대시보드'에서 경로·격리·설정·체크리스트·감사이력을 다루는 "
                  "하나의 운영 콘솔로 진화했습니다.", top=6.1, height=0.75)

    # 12 · GOVERNANCE — physical 20
    s = body_slide(prs, blank, "12 — 거버넌스", 20, TOTAL, "12 · GOVERNANCE",
                   "관리 기능 고도화 — 누가, 무엇을 바꿨는가")
    bullets(s, [
        "—  작업자 식별 — 이름/사번 1회 입력, 이후 모든 변경에 자동 첨부",
        '—  관리자 변경 이력 — 경로·격리·설정·체크리스트 변경을 "누가·언제·무엇을" 시간순 기록',
        "—  설치 체크리스트 — 신규 구축 8단계를 완료자·완료 시각과 함께 추적",
        "—  설정 자동 백업 — 변경마다 스냅샷 보관(최근 20개), 실수 복구 가능",
    ], top=2.0, size=13, line_h=0.8)
    impact_box(s, "'누가 바꿨는지 알 수 없다'는 감사의 공백을 관리자 조작까지 포함해 없앴습니다. "
                  "여러 개발자가 나눠 구축해도 진행 상황이 한눈에 보입니다.", top=5.55, height=1.15)

    # 13 · SUMMARY & SI-READY — physical 21
    s = body_slide(prs, blank, "13 — 요약·SI표준", 21, TOTAL, "13 · SUMMARY & SI-READY",
                   "비전문가용 요약 & SI/온프레미스 표준화")
    bullets(s, [
        '—  요약 탭 — 상태 배너("정상 보호"/"확인 필요"/"응답 없음") + 문장형 이벤트로 즉시 파악',
        "—  LHAMS_HOME 한 줄 설정 — 경로 7개 대신 홈 경로 하나로 구축 시간 단축",
        "—  사이트 식별(SITE_NAME/ID) + /api/health — 다중 서버 가동상태 한눈에 확인",
        "—  폐쇄망 대응 — 런타임 외부 호출 없음, 빌드 산출물만 반입하는 오프라인 배포",
    ], top=2.0, size=13, line_h=0.8)
    impact_box(s, "세미나용 PoC에서 SI·온프레미스 다중 고객사에 표준화된 방식으로 구축·운영할 수 있는 "
                  "제품 형태로 발전했습니다.", top=5.55, height=1.15)

    # ---- 3) 슬라이드 순서 재배치: 새 4장을 물리적 18번째 위치(구 10 레이아웃 슬라이드 뒤)로 이동 ----
    xml_slides = slides._sldIdLst
    slide_ids = list(xml_slides)
    new_ids = slide_ids[-4:]           # 방금 추가한 4장 (끝에 붙어있음)
    for sid in new_ids:
        xml_slides.remove(sid)
    insert_at = 17                      # 0-indexed: 기존 index16(물리 17번, 프로젝트 레이아웃) 바로 다음
    for offset, sid in enumerate(new_ids):
        xml_slides.insert(insert_at + offset, sid)

    # ---- 4) 전체 재배치 후 헤더/섹션/페이지 번호 최종 정리 ----
    # 새 슬라이드 4장은 생성 시 이미 올바른 라벨/페이지 번호로 만들어졌으므로 손댈 필요 없음.
    # 뒤로 밀린 기존 슬라이드(옛 물리 18·19·20 → 새 물리 22·23·24)만 라벨을 갱신한다.
    demo_slide = slides[21]   # 0-indexed 21 == 새 물리 22번째 슬라이드 (11 라이브 데모)
    update_header_footer(demo_slide, top_right_label="15 — 라이브 데모", tag="15 · LIVE DEMO", page_no=22, total=TOTAL)

    conclusion_slide = slides[22]  # 새 물리 23번째 (12 결론)
    update_header_footer(conclusion_slide, top_right_label="16 — 결론", tag="16 · CONCLUSION", page_no=23, total=TOTAL)

    closing_slide = slides[23]  # 새 물리 24번째 (THANK YOU)
    update_header_footer(closing_slide, page_no=24, total=TOTAL)

    # 나머지 기존 바디 슬라이드(물리 2~17)는 위치가 그대로이므로 "NN / 20" → "NN / 24"만 갱신
    for idx in range(1, 17):
        try:
            update_header_footer(slides[idx], page_no=idx + 1, total=TOTAL)
        except KeyError:
            pass  # 커버 슬라이드 등 헤더가 없는 레이아웃은 건너뜀

    prs.save(PPTX_PATH)
    print(f"OK: saved {PPTX_PATH} ({len(prs.slides._sldIdLst)} slides, expected 24)")


if __name__ == "__main__":
    main()
